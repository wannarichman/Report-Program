# -*- coding: utf-8 -*-
"""Report.py — AI Live Sync Master Builder (v4 전체 정상본)
원본 v1의 모든 기능 보존 + v3 개선(A/B/C1/C2)
"""
import streamlit as st
import streamlit.components.v1 as components
import pandas as pd
import json
import time
import base64
import os
import copy
import urllib.parse
import requests
import google.generativeai as genai
from datetime import datetime
import zoneinfo
import re
import struct
import zipfile
import xml.etree.ElementTree as ET
from io import BytesIO

try:
    from PIL import Image
    HAS_PIL = True
except Exception:
    HAS_PIL = False
    
from pathlib import Path

try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except Exception:
    HAS_PYPDF = False

MAX_DOC_CHARS = 30000
COMPANY_FACTS_PATH = "posco_projects.json"
SQ = chr(39)  # 작은따옴표. f-string 내 이스케이프 회피용
ERROR_IMG = "https://placehold.co/800x400/f8fafc/94a3b8?text=Image+Not+Found"


# ==========================================
# 0. 유틸
# ==========================================
def _strip_code_fence(text):
    t = (text or "").strip()
    fence = "`" * 3
    if t.startswith(fence):
        nl = t.find("\n")
        if nl != -1:
            t = t[nl + 1:]
        if t.endswith(fence):
            t = t[:-3]
    return t.strip()


def extract_requested_page_count(text):
    if not text:
        return None
    m = re.search(r"(\d+)\s*(페이지|장|쪽|p)", str(text))
    if m:
        try:
            n = int(m.group(1))
            if 1 <= n <= 30:
                return n
        except Exception:
            pass
    return None


# ==========================================
# 0-1. 포맷별 텍스트 추출
# ==========================================
def extract_pdf(raw):
    if not HAS_PYPDF:
        return "[오류] pypdf 미설치"
    try:
        reader = PdfReader(BytesIO(raw))
        out = []
        for i, page in enumerate(reader.pages):
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            if t.strip():
                out.append(f"--- [PDF p.{i+1}] ---\n{t.strip()}")
        full = "\n\n".join(out).strip()
        return full or "[알림] PDF에서 텍스트를 추출하지 못했습니다(스캔본/이미지 PDF 추정)."
    except Exception as e:
        return f"[PDF 파싱 오류] {e}"


def extract_pptx(raw):
    try:
        from pptx import Presentation
    except ImportError:
        return "[오류] python-pptx 미설치"
    try:
        prs = Presentation(BytesIO(raw))
        out = []
        for i, slide in enumerate(prs.slides):
            buf = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            buf.append(line)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        joined = " | ".join([c for c in cells if c])
                        if joined:
                            buf.append(joined)
            try:
                if slide.has_notes_slide:
                    note = slide.notes_slide.notes_text_frame.text.strip()
                    if note:
                        buf.append(f"[노트] {note}")
            except Exception:
                pass
            if buf:
                out.append(f"--- [Slide {i+1}] ---\n" + "\n".join(buf))
        return "\n\n".join(out).strip() or "[알림] PPTX에서 텍스트를 찾지 못했습니다."
    except Exception as e:
        return f"[PPTX 파싱 오류] {e}"


def extract_docx(raw):
    try:
        from docx import Document
    except ImportError:
        return "[오류] python-docx 미설치"
    try:
        doc = Document(BytesIO(raw))
        out = []
        for para in doc.paragraphs:
            if para.text.strip():
                out.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                line = " | ".join(c.text.strip() for c in row.cells)
                if line.strip():
                    out.append(line)
        return "\n".join(out).strip() or "[알림] DOCX에서 텍스트를 찾지 못했습니다."
    except Exception as e:
        return f"[DOCX 파싱 오류] {e}"


def extract_hwpx(raw):
    try:
        out = []
        with zipfile.ZipFile(BytesIO(raw)) as zf:
            names = sorted([
                n for n in zf.namelist()
                if n.startswith("Contents/section") and n.endswith(".xml")
            ])
            for n in names:
                try:
                    root = ET.fromstring(zf.read(n))
                except Exception:
                    continue
                texts = []
                for elem in root.iter():
                    tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                    if tag == "t" and elem.text:
                        texts.append(elem.text)
                if texts:
                    out.append("\n".join(texts))
        return "\n\n".join(out).strip() or "[알림] HWPX에서 텍스트를 찾지 못했습니다."
    except Exception as e:
        return f"[HWPX 파싱 오류] {e}"


def extract_hwp(raw):
    try:
        import olefile
        import zlib
    except ImportError:
        return "[오류] olefile 미설치"
    try:
        ole = olefile.OleFileIO(BytesIO(raw))
        is_compressed = True
        try:
            header = ole.openstream("FileHeader").read()
            if len(header) > 36:
                is_compressed = bool(header[36] & 0x01)
        except Exception:
            pass

        sections = [
            s for s in ole.listdir(streams=True)
            if "/".join(s).startswith("BodyText/Section")
        ]
        if not sections:
            ole.close()
            return "[알림] HWP에서 BodyText 섹션을 찾지 못했습니다(암호화/구버전 가능성). HWPX로 저장 후 다시 시도해 보세요."

        chunks = []
        for stream in sections:
            try:
                raw_stream = ole.openstream(stream).read()
                data = zlib.decompress(raw_stream, -15) if is_compressed else raw_stream
                i = 0
                paras = []
                while i + 4 <= len(data):
                    head = struct.unpack("<I", data[i:i+4])[0]
                    tag_id = head & 0x3FF
                    size = (head >> 20) & 0xFFF
                    i += 4
                    if size == 0xFFF:
                        if i + 4 > len(data):
                            break
                        size = struct.unpack("<I", data[i:i+4])[0]
                        i += 4
                    payload = data[i:i+size]
                    i += size
                    if tag_id == 67:  # HWPTAG_PARA_TEXT
                        try:
                            text = payload.decode("utf-16-le", errors="ignore")
                            cleaned = "".join(
                                ch for ch in text
                                if ord(ch) >= 0x20 or ch in "\n\t"
                            )
                            if cleaned.strip():
                                paras.append(cleaned)
                        except Exception:
                            pass
                if paras:
                    chunks.append("\n".join(paras))
            except Exception:
                continue
        ole.close()
        return "\n\n".join(chunks).strip() or "[알림] HWP 본문 텍스트를 추출하지 못했습니다."
    except Exception as e:
        return f"[HWP 파싱 오류] {e}"


def extract_text_from_upload(uploaded_file):
    if uploaded_file is None:
        return ""
    name = (uploaded_file.name or "").lower()
    raw = uploaded_file.getvalue()
    if name.endswith(".pdf"):
        text = extract_pdf(raw)
    elif name.endswith(".pptx"):
        text = extract_pptx(raw)
    elif name.endswith(".ppt"):
        text = "[알림] .ppt(구버전) 형식은 직접 파싱이 제한됩니다. PowerPoint에서 '다른 이름으로 저장 → .pptx' 후 다시 업로드해 주세요."
    elif name.endswith(".docx"):
        text = extract_docx(raw)
    elif name.endswith(".hwpx"):
        text = extract_hwpx(raw)
    elif name.endswith(".hwp"):
        text = extract_hwp(raw)
    else:
        try:
            text = raw.decode("utf-8", errors="ignore")
        except Exception:
            text = str(raw)
    return (text or "")[:MAX_DOC_CHARS]


def load_uploaded_images(files):
    """Streamlit UploadedFile 리스트 → Gemini Vision 입력·토큰 치환에 쓸 dict 리스트.
    각 항목: { index, name, mime, data_url, raw, pil }
    """
    out = []
    if not files:
        return out
    for idx, f in enumerate(files):
        try:
            raw = f.getvalue()
            if not raw:
                continue
            mime = (getattr(f, "type", "") or "").lower()
            if not mime.startswith("image/"):
                ext = (f.name or "").lower().split(".")[-1]
                mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext or 'jpeg'}"
            b64 = base64.b64encode(raw).decode()
            entry = {
                "index": idx,
                "name": f.name or f"photo_{idx}",
                "mime": mime,
                "data_url": f"data:{mime};base64,{b64}",
                "raw": raw,
                "pil": None,
            }
            if HAS_PIL:
                try:
                    entry["pil"] = Image.open(BytesIO(raw))
                except Exception:
                    entry["pil"] = None
            out.append(entry)
        except Exception:
            continue
    return out


def resolve_uploaded_tokens(report, images):
    """AI가 생성한 JSON 안의 'UPLOADED:N' 토큰을 실제 base64 data URL로 치환."""
    if not images or not isinstance(report, dict):
        return report
    url_map = {f"UPLOADED:{im['index']}": im["data_url"] for im in images}

    def fix(v):
        if isinstance(v, str):
            s = v.strip()
            if s.startswith("UPLOADED:"):
                return url_map.get(s, v)
        return v

    for pg in report.get("pages", []) or []:
        for sec in pg.get("sections", []) or []:
            sec["main_image"] = fix(sec.get("main_image"))
            for it in sec.get("side_items", []) or []:
                if it.get("type") == "image":
                    it["src"] = fix(it.get("src"))
    return report

# ==========================================
# 1. 페이지 설정 및 CSS
# ==========================================
st.set_page_config(page_title="AI Live Sync Master Builder", layout="wide")

st.markdown("""
<style>
[data-testid="stAppViewContainer"] { background-color: #ffffff !important; }
.main [data-testid="stVerticalBlockBorderWrapper"] {
    background-color: #ffffff !important; border: 1px solid #dee2e6 !important;
    border-radius: 16px !important; padding: 35px 40px !important;
    box-shadow: 0 4px 16px rgba(0,0,0,0.04) !important; margin-bottom: 50px !important;
}
[data-testid="stSidebar"] [data-testid="stVerticalBlockBorderWrapper"] {
    border: 1px solid #dee2e6 !important; padding: 15px !important;
    box-shadow: none !important; margin-bottom: 10px !important;
}
.side-slot-card { padding: 10px 0px; margin-bottom: 16px; }
.text-line { white-space: pre-wrap; word-wrap: break-word; line-height: 1.8; margin-bottom: 10px; color: #334155; }
.voice-panel { background:#fff; border:1px solid #dee2e6; padding:15px; border-radius:16px; text-align:center; margin-bottom:15px; }
.btn-mute { padding:8px 16px; background:#6c757d; color:#fff; border:none; border-radius:8px; cursor:pointer; font-weight:bold; width:100%; }
.btn-mute.active { background:#dc3545; }
</style>
""", unsafe_allow_html=True)


# ==========================================
# 2. 전역 저장소
# ==========================================
@st.cache_resource
def get_global_store():
    return {
        "report_data": None, "current_page": 0, "user_labels": {},
        "chat_history": [], "active_sessions": {},
        "voice_channel": "posco_briefing_room",
    }

shared_store = get_global_store()


# ==========================================
# 3. 표준 양식
# ==========================================
def get_sample_json_guide():
    return {
        "title": "주간 보고 (AI 생성 기반)", "title_fs": 55, "title_color": "#0f172a",
        "pages": [
            {"tab": "요약", "header": "Executive Summary", "header_fs": 35, "header_color": "#475569",
             "sections": [{"title": "핵심 요약", "title_fs": 32, "title_color": "#1a1c1e", "col_ratio": 1.5,
                           "main_image": None, "full_width": True, "image_query": "",
                           "chart_type": "Bar", "chart_data": "",
                           "lines": [{"text": "• 금주 핵심 성과를 요약합니다.", "size": 24, "color": "#1e293b"},
                                     {"text": "• 주요 이슈 및 리스크를 점검합니다.", "size": 22, "color": "#1e293b"}],
                           "side_items": [{"type": "metric", "label": "종합 진행률", "value": "0%", "color": "#007bff", "label_fs": 14, "label_color": "#64748b", "value_fs": 34}]}]},
            {"tab": "상세 (데이터)", "header": "현황 상세 및 데이터 분석", "header_fs": 35, "header_color": "#475569",
             "sections": [{"title": "데이터 지표 분석", "title_fs": 32, "title_color": "#1a1c1e", "col_ratio": 1.5,
                           "main_image": None, "full_width": True, "image_query": "",
                           "chart_type": "Bar", "chart_data": "1분기, 35\n2분기, 50\n3분기, 42\n4분기, 68",
                           "lines": [{"text": "• 위 데이터 차트를 통해 실적 추이를 확인할 수 있습니다.", "size": 22, "color": "#1e293b"}],
                           "side_items": [{"type": "metric", "label": "정량 지표 요약", "value": "목표 달성률: 85%", "color": "#16a34a", "label_fs": 14, "label_color": "#64748b", "value_fs": 22}]}]},
            {"tab": "액션/리스크", "header": "Action Items & Risks", "header_fs": 35, "header_color": "#475569",
             "sections": [{"title": "Action Items", "title_fs": 32, "title_color": "#1a1c1e", "col_ratio": 1.5,
                           "main_image": None, "full_width": True, "image_query": "",
                           "chart_type": "Bar", "chart_data": "",
                           "lines": [{"text": "1) [담당/기한] 해결 과제 1", "size": 22, "color": "#1e293b"}],
                           "side_items": [{"type": "metric", "label": "주요 Blocker", "value": "없음", "color": "#dc2626", "label_fs": 14, "label_color": "#64748b", "value_fs": 26}]}]},
        ],
    }


def create_empty_page():
    return copy.deepcopy(get_sample_json_guide()["pages"][0])


# ==========================================
# 3-1. JSON 자동 보정
# ==========================================
DEFAULT_LINE = {"text": "", "size": 22, "color": "#1e293b"}
DEFAULT_METRIC = {
    "type": "metric", "label": "항목", "value": "",
    "color": "#007bff", "label_fs": 14, "label_color": "#64748b", "value_fs": 28,
}
DEFAULT_IMAGE_ITEM = {"type": "image", "src": None, "width": 350, "image_query": ""}


def _normalize_section(sec):
    sec.setdefault("title", "섹션")
    sec.setdefault("title_fs", 32)
    sec.setdefault("title_color", "#1a1c1e")
    sec.setdefault("col_ratio", 1.5)
    sec.setdefault("main_image", None)
    sec.setdefault("full_width", True)
    sec.setdefault("image_query", "")
    sec.setdefault("chart_type", "Bar")
    sec.setdefault("chart_data", "")

    raw_lines = sec.get("lines") or []
    fixed_lines = []
    for ln in raw_lines:
        if isinstance(ln, str):
            fixed_lines.append({**DEFAULT_LINE, "text": ln})
        elif isinstance(ln, dict):
            base = dict(DEFAULT_LINE)
            base.update({k: v for k, v in ln.items() if v is not None})
            fixed_lines.append(base)
    sec["lines"] = fixed_lines

    raw_items = sec.get("side_items") or []
    fixed_items = []
    for it in raw_items:
        if not isinstance(it, dict):
            continue
        t = it.get("type", "metric")
        if t == "metric":
            base = dict(DEFAULT_METRIC)
            base.update({k: v for k, v in it.items() if v is not None})
            fixed_items.append(base)
        elif t == "image":
            base = dict(DEFAULT_IMAGE_ITEM)
            base.update({k: v for k, v in it.items() if v is not None})
            fixed_items.append(base)
    sec["side_items"] = fixed_items
    return sec


def adapt_json_format(raw_data):
    if not isinstance(raw_data, dict) or "pages" not in raw_data or not isinstance(raw_data["pages"], list):
        return get_sample_json_guide()
    raw_data.setdefault("title", "AI 자동 생성 보고서")
    raw_data.setdefault("title_fs", 55)
    raw_data.setdefault("title_color", "#0f172a")
    fixed_pages = []
    for pg in raw_data["pages"]:
        if not isinstance(pg, dict):
            continue
        pg.setdefault("tab", "페이지")
        pg.setdefault("header", "")
        pg.setdefault("header_fs", 35)
        pg.setdefault("header_color", "#475569")
        secs = pg.get("sections") or []
        pg["sections"] = [_normalize_section(s) for s in secs if isinstance(s, dict)]
        if not pg["sections"]:
            pg["sections"] = [_normalize_section({"title": pg["header"] or pg["tab"], "lines": []})]
        fixed_pages.append(pg)
    if not fixed_pages:
        return get_sample_json_guide()
    raw_data["pages"] = fixed_pages
    return raw_data


# ==========================================
# 4. 외부 정보 수집 (네이버 검색)  ★ v3 (A) 강화
# ==========================================
def naver_search_text(query, max_results=5):
    """뉴스(최신순) + 백과 + 웹 멀티 소스 통합 그라운딩."""
    try:
        cid = st.secrets.get("NAVER_CLIENT_ID", "")
        csec = st.secrets.get("NAVER_CLIENT_SECRET", "")
        if not cid or not csec:
            return ""
        headers = {"X-Naver-Client-Id": cid, "X-Naver-Client-Secret": csec}
        snippets = []
        endpoints = [
            ("news.json", {"sort": "date"}),
            ("encyc.json", {"sort": "sim"}),
            ("webkr.json", {"sort": "sim"}),
        ]
        for endpoint, extra in endpoints:
            try:
                params = {"query": query, "display": max_results}
                params.update(extra)
                r = requests.get(
                    "https://openapi.naver.com/v1/search/" + endpoint,
                    params=params, headers=headers, timeout=5,
                )
                if r.status_code == 200:
                    for item in r.json().get("items", []):
                        title = re.sub(r"<[^>]+>", "", item.get("title", ""))
                        desc = re.sub(r"<[^>]+>", "", item.get("description", ""))
                        link = item.get("link", "")
                        pub = item.get("pubDate", "")
                        tag = endpoint.split(".")[0]
                        snippets.append(f"- [{tag}/{pub}] {title}: {desc} (출처: {link})")
            except Exception:
                continue
        return "\n".join(snippets[:20])
    except Exception:
        return ""


def naver_search_image(query):
    """★ v3 (C2): query 핵심어가 결과 제목에 포함된 이미지를 우선 채택. 없으면 1순위 폴백."""
    try:
        cid = st.secrets.get("NAVER_CLIENT_ID", "")
        csec = st.secrets.get("NAVER_CLIENT_SECRET", "")
        if not cid or not csec:
            return ""
        r = requests.get(
            "https://openapi.naver.com/v1/search/image",
            params={"query": query, "display": 10, "sort": "sim", "filter": "large"},
            headers={"X-Naver-Client-Id": cid, "X-Naver-Client-Secret": csec},
            timeout=5,
        )
        if r.status_code != 200:
            return ""
        items = r.json().get("items", [])
        if not items:
            return ""
        keywords = [w for w in re.split(r"\s+", query.strip()) if len(w) >= 2]
        for it in items:
            title = re.sub(r"<[^>]+>", "", it.get("title", "")).lower()
            if any(k.lower() in title for k in keywords):
                return it.get("link", "")
        return items[0].get("link", "")
    except Exception:
        return ""


def get_auto_image_url(query, w=1600, h=900):
    if not isinstance(query, str):
        return ""
    q = query.strip()
    if not q:
        return ""
    naver_url = naver_search_image(q)
    if naver_url:
        return naver_url
    tags = ",".join([urllib.parse.quote(t.strip()) for t in q.split(",") if t.strip()])
    lock = abs(hash(q)) % 100000
    return f"https://loremflickr.com/{w}/{h}/{tags}?lock={lock}"


def render_image_src(img_val):
    if not img_val or not isinstance(img_val, str):
        return ""
    val = img_val.strip()
    if not val:
        return ""
    if val.startswith("http://") or val.startswith("https://") or val.startswith("data:image"):
        return val
    if os.path.isfile(val):
        try:
            with open(val, "rb") as f:
                ext = val.split(".")[-1].lower()
                mime = "image/jpeg" if ext in ["jpg", "jpeg"] else f"image/{ext}"
                return f"data:{mime};base64,{base64.b64encode(f.read()).decode()}"
        except Exception:
            return val
    return val


# ==========================================
# 4-1. 회사 시공현황 사실 DB  ★ v3 (B) 신규
# ==========================================
DEFAULT_COMPANY_FACTS = [
    {"name": "청라-영종 하늘대교", "type": "해상교량", "location": "인천 청라~영종", "status": "시공중", "completion": "2027", "scale": "총연장 약 4.05km, 사장교 구간 포함"},
    {"name": "새만금 남북도로 2단계", "type": "도로", "location": "전북 새만금", "status": "시공중"},
    {"name": "송도 R3 블록", "type": "공동주택", "location": "인천 송도국제도시", "status": "분양/시공"},
    {"name": "광양 LNG 터미널", "type": "플랜트/에너지", "location": "전남 광양", "status": "시공중"},
    {"name": "포항 이차전지 소재 공장", "type": "산업플랜트", "location": "경북 포항", "status": "시공중"},
]


def load_company_facts():
    p = Path(COMPANY_FACTS_PATH)
    if p.exists():
        try:
            return json.loads(p.read_text(encoding="utf-8"))
        except Exception:
            return DEFAULT_COMPANY_FACTS
    return DEFAULT_COMPANY_FACTS


def format_facts_for_prompt(facts):
    lines = []
    for f in facts:
        parts = [f"- {f.get('name', '')}"]
        for k in ("type", "location", "status", "completion", "scale"):
            if f.get(k):
                parts.append(f"{k}={f[k]}")
        lines.append(" | ".join(parts))
    return "\n".join(lines)


# ==========================================
# 5. AI 텍스트 -> JSON  ★ v3 (B + C1) 강화
# ==========================================
def generate_json_from_ai(api_key, context_text, requested_pages=None, images=None):
    try:
        genai.configure(api_key=api_key)
        try:
            available = []
            for m in genai.list_models():
                methods = getattr(m, "supported_generation_methods", []) or []
                if "generateContent" in methods:
                    available.append(m.name)
        except Exception as e:
            return {"error": f"모델 목록 조회 실패: {e}"}
        if not available:
            return {"error": "이 API 키로 generateContent 가능한 모델이 없습니다."}

        preferred_order = [
            "gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.0-flash-001",
            "gemini-1.5-flash", "gemini-1.5-flash-latest", "gemini-1.5-flash-002",
            "gemini-1.5-pro", "gemini-1.5-pro-latest",
        ]

        def pick_model(avail_list):
            short_names = {name.split("/")[-1]: name for name in avail_list}
            for p in preferred_order:
                if p in short_names:
                    return short_names[p]
            for n in avail_list:
                if "flash" in n:
                    return n
            return avail_list[0]

        chosen_model = pick_model(available)

        now_kst = datetime.now(zoneinfo.ZoneInfo("Asia/Seoul"))
        today_str = now_kst.strftime("%Y년 %m월 %d일")
        year_str = now_kst.strftime("%Y")
        quarter = (now_kst.month - 1) // 3 + 1

        # ★ (A) 강화: 두 쿼리로 폭넓게 그라운딩
        search_seed = (context_text or "")[:200].replace("\n", " ").strip()
        web_context_parts = []
        if search_seed:
            web_context_parts.append(naver_search_text(search_seed))
            web_context_parts.append(naver_search_text("포스코이앤씨 " + search_seed[:80]))
        web_context = "\n".join([p for p in web_context_parts if p]).strip()
        if not web_context:
            web_context = "(외부 검색 결과 없음 - 입력 데이터 기반으로만 생성)"

        # ★ (B) 신규: 회사 시공현황 사실 DB
        company_facts = format_facts_for_prompt(load_company_facts())

        page_count_directive = (
            f"\n[페이지 수 제약]\n- pages 배열은 정확히 {requested_pages}개로 만들 것.\n"
            if requested_pages else ""
        )

        schema_doc = json.dumps(get_sample_json_guide(), ensure_ascii=False, indent=2)

        # ★ (C1) 강화: image_query 규칙을 고유명사+핵심어 강제로 변경
        system_prompt = (
            f"당신은 포스코이앤씨에서 쓰이는 주간/프로젝트 보고서 JSON 생성기입니다.\n"
            "입력 데이터를 분석하여 보고서용 JSON을 생성하세요.\n\n"
            f"[현재 시점]\n- 오늘: {today_str}\n- 연도: {year_str}년 / 분기: {year_str}년 {quarter}분기\n"
            "- 과거 연도(2024, 2025년)를 '현재'로 표현 금지.\n\n"
            "[출력 형식]\n- 순수 JSON 하나만 출력 (마크다운/코드펜스 절대 금지).\n"
            "- 최상위 키: title, title_fs, title_color, pages.\n"
            "- 각 페이지는 sections 최소 1개. 각 섹션은 lines 최소 2개, side_items 최소 1개.\n"
            f"{page_count_directive}"
            f"[JSON 스키마 - 필드 이름만 참고]\n{schema_doc}\n\n"
            f"[회사 시공현황 사실 DB]  ★ 본문/지표/일정/위치 등을 인용할 때 우선 사용\n{company_facts}\n\n"
            f"[외부 검색 컨텍스트]  ★ 네이버 뉴스/백과/웹 실시간 결과\n{web_context}\n\n"
            "[절대 규칙]\n"
            "1. 환각 금지: 팀명/부서명/인물/프로젝트명/회사명/수치/날짜는 [회사 시공현황 사실 DB] 또는 [외부 검색 컨텍스트] 또는 [입력 데이터]에 명시된 것만 사용.\n"
            "2. 위 자료에 없는 사실은 비우거나(\"\") 일반적 표현(\"관련 부서\", \"담당자\") 사용. 절대 일반론·추측·창작 금지.\n"
            "3. chart_data는 실제 수치 있을 때만 '항목, 수치' 줄 형태로 개행구분. 없으면 완전히 빈 문자열.\n"
            "4. image_query 작성 규칙 (중요):\n"
            "   - 프로젝트 고유명사 + 핵심 키워드 결합. 예) '청라 하늘대교 사장교 시공', '새만금 남북도로', 'Cheongna Sky Bridge construction'.\n"
            "   - 일반 명사 단독 금지: '교량', '건설현장', '공사', 'report', 'summary', 'data' 등 ❌.\n"
            "   - 한국 프로젝트/지명은 한글, 글로벌 키워드는 영문도 허용.\n"
            "   - 단서가 없으면 빈 문자열.\n"
            "5. 표준양식의 플레이스홀더 문구를 그대로 복사 금지. 입력 데이터 + 검색 컨텍스트 + 회사 사실 DB 기반으로 재작성.\n\n"
            f"[입력 데이터]\n{context_text}\n"
        )

        # ★ v6: 첨부 사진 규칙을 시스템 프롬프트에 주입
        n_images = len(images) if images else 0
        if n_images > 0:
            image_directive = (
                f"\n[첨부 현장 사진]  ★ 총 {n_images}장 (인덱스 0..{n_images - 1})\n"
                "- 각 사진을 직접 분석해서 본문(lines)에 최대한 구체적으로 설명하세요 (구조물·인원·장비·안전이슈·계절·날씨 등).\n"
                "- 사진을 보고서 안에 실제로 배치하려면 main_image 또는 side_items[].src 필드에 "
                "\"UPLOADED:0\", \"UPLOADED:1\" 형식의 토큰을 적으세요 (외부 URL·더미 이미지 금지).\n"
                "- 한 사진을 여러 섹션에 재사용 가능. 관련성 없는 사진은 생략해도 됨.\n"
                "- 가능하면 첫 섹션 main_image에 가장 대표적인 사진(0번 우선) 배치.\n"
            )
            system_prompt = system_prompt + image_directive

        # 멀티모달 parts 구성: [text, img1, img2, ...]
        parts = [system_prompt]
        if images:
            for im in images:
                if im.get("pil") is not None:
                    parts.append(im["pil"])
                else:
                    try:
                        parts.append({"mime_type": im["mime"], "data": im["raw"]})
                    except Exception:
                        continue

        generation_config = {"response_mime_type": "application/json", "temperature": 0.5}
        tried = []
        candidates = [chosen_model] + [n for n in available if n != chosen_model]
        response = None
        for model_name in candidates:
            try:
                model = genai.GenerativeModel(model_name, generation_config=generation_config)
                response = model.generate_content(parts if len(parts) > 1 else parts[0])
                break
            except Exception as e:
                tried.append(f"{model_name}: {e}")
                continue
                
        if response is None:
            return {"error": f"모든 모델 호출 실패: {tried}"}

        clean_text = _strip_code_fence(response.text or "")
        try:
            parsed = json.loads(clean_text)
        except Exception as e:
            return {"error": f"JSON 파싱 실패: {e}\n동봉(앞 500자): {clean_text[:500]}"}

        n_pages = len(parsed.get("pages", []))
        n_lines_total = sum(
            len(s.get("lines", []))
            for p in parsed.get("pages", [])
            for s in p.get("sections", [])
        )
        if n_pages == 0 or n_lines_total == 0:
            return {"error": "모델이 빈 보고서를 돌려주었습니다. 입력을 더 자세히 적어주세요."}
        return parsed

    except Exception as e:
        return {"error": str(e)}

# ==========================================
# 5-1. 입장 게이트 (역할 선택 + 이름 입력)
# ==========================================
def render_onboarding_gate():
    """역할/이름이 정해지지 않았으면 게이트 화면을 띄우고 False 반환.
    True가 반환된 경우에만 메인 UI를 그려야 한다.
    """
    if st.session_state.get("user_ready"):
        return True

    st.markdown(
        "<div style='text-align:center; padding:30px 0 10px 0;'>"
        "<h1 style='margin:0; color:#0f172a;'>📋 AI Live Sync 입장</h1>"
        "<p style='color:#64748b; margin-top:8px;'>역할을 선택하고 이름(또는 직함)을 입력해 주세요.</p>"
        "</div>",
        unsafe_allow_html=True,
    )

    col_l, col_c, col_r = st.columns([1, 2, 1])
    with col_c:
        with st.container(border=True):
            with st.form("onboarding_form", clear_on_submit=False):
                role_label = st.radio(
                    "역할",
                    ["🎤 보고자 (발표·편집 권한)", "👥 피보고자 (청취·채팅)"],
                    index=1,
                    horizontal=False,
                )
                name = st.text_input(
                    "이름 / 직함",
                    value=st.session_state.get("user_name", ""),
                    placeholder="예) 최우석 대리 / 안전팀 김OO",
                    max_chars=30,
                )
                submitted = st.form_submit_button("입장", use_container_width=True, type="primary")

            if submitted:
                if not (name and name.strip()):
                    st.error("이름을 입력해 주세요.")
                    return False
                role_key = "reporter" if role_label.startswith("🎤") else "audience"
                role_kor = "보고자" if role_key == "reporter" else "피보고자"
                st.session_state.user_role = role_key
                st.session_state.user_name = name.strip()
                st.session_state.user_label = f"{name.strip()} ({role_kor})"
                st.session_state.user_ready = True
                # 새로고침 시에도 유지되도록 URL 파라미터에 보존
                st.query_params["role"] = role_key
                st.query_params["name"] = name.strip()
                st.rerun()

    st.caption("💡 같은 링크를 여러 명이 열어도 각자 입력한 이름으로 구분됩니다.")
    return False
    
# ==========================================
# 6. ID 식별 및 Agora 음성
# ==========================================
# ==========================================
# 6. ID 식별 (uid 발급 + 역할/이름 복원)
# ==========================================
if "uid" not in st.session_state:
    url_uid = st.query_params.get("uid")
    if url_uid:
        st.session_state.uid = url_uid
    else:
        new_uid = f"u{int(time.time() * 1000)}"
        st.session_state.uid = new_uid
        st.query_params["uid"] = new_uid

# query_params에 role/name이 있으면 자동 복원 (새로고침 시 게이트 건너뜀)
if "user_ready" not in st.session_state:
    qr_role = st.query_params.get("role")
    qr_name = st.query_params.get("name")
    if qr_role in ("reporter", "audience") and qr_name:
        role_kor = "보고자" if qr_role == "reporter" else "피보고자"
        st.session_state.user_role = qr_role
        st.session_state.user_name = qr_name
        st.session_state.user_label = f"{qr_name} ({role_kor})"
        st.session_state.user_ready = True

# 게이트 통과 못하면 사이드바/메인 렌더 중단
if not render_onboarding_gate():
    st.stop()

# 이 시점부터는 user_role / user_name / user_label 모두 보장됨
is_reporter = (st.session_state.user_role == "reporter")
my_label = st.session_state.user_label

def agora_voice_system(app_id, channel, user_label):
    """Agora 음성:
    - 연결 상태 (CONNECTING / CONNECTED / DISCONNECTED / FAILED) 를 시각적으로 표시
    - 마이크 음소거(mute) 와 채널 연결(connect) 을 분리
    - 연결 끊김 시 최대 3회 자동 재연결
    - 음량 레벨바는 '연결됨' & '음소거 해제' 일 때만 동작
    """
    custom_html = """
<script src=\"https://download.agora.io/sdk/release/AgoraRTC_N-4.20.0.js\"></script>
<style>
.voice-status { display:flex; align-items:center; gap:8px; justify-content:center; font-size:13px; font-weight:700; margin-bottom:8px; }
.voice-status .dot { width:10px; height:10px; border-radius:50%; display:inline-block; }
.voice-status.connecting .dot { background:#f59e0b; animation:pulse 1s infinite; }
.voice-status.connected  .dot { background:#22c55e; box-shadow:0 0 0 4px rgba(34,197,94,0.18); }
.voice-status.failed     .dot { background:#ef4444; }
.voice-status.disconnected .dot { background:#94a3b8; }
.voice-status.connecting    { color:#92400e; }
.voice-status.connected     { color:#166534; }
.voice-status.failed        { color:#b91c1c; }
.voice-status.disconnected  { color:#475569; }
@keyframes pulse { 0%,100% { opacity:1; } 50% { opacity:0.35; } }
.btn-mute { padding:8px 16px; background:#22c55e; color:#fff; border:none; border-radius:8px; cursor:pointer; font-weight:bold; width:100%; }
.btn-mute.muted { background:#ef4444; }
.btn-mute[disabled] { background:#cbd5e1; cursor:not-allowed; }
</style>

<div class=\"voice-panel\">
    <div id=\"v-status\" class=\"voice-status connecting\"><span class=\"dot\"></span><span id=\"v-status-text\">연결 중…</span></div>
    <div style=\"font-size:11px; text-align:center; color:#64748b; margin-bottom:8px;\">USER_LABEL</div>
    <div style=\"width:100%; height:10px; background:#e2e8f0; border-radius:5px; margin-bottom:12px; overflow:hidden;\">
        <div id=\"level-bar\" style=\"width:0%; height:100%; background:#22c55e; transition:width 0.05s;\"></div>
    </div>
    <button id=\"mute\" class=\"btn-mute\" disabled>🎤 마이크 : 켜짐</button>
</div>

<script>
let client = AgoraRTC.createClient({ mode: \"rtc\", codec: \"vp8\" });
let localTracks = { audioTrack: null };
let isMuted = false;
let retryCount = 0;
const MAX_RETRY = 3;

function setStatus(state, text) {
    const el = document.getElementById(\"v-status\");
    if (!el) return;
    el.classList.remove(\"connecting\", \"connected\", \"failed\", \"disconnected\");
    el.classList.add(state);
    document.getElementById(\"v-status-text\").innerText = text;
    const btn = document.getElementById(\"mute\");
    if (btn) btn.disabled = (state !== \"connected\");
    if (state !== \"connected\") {
        const lb = document.getElementById(\"level-bar\");
        if (lb) lb.style.width = \"0%\";
    }
}

async function join() {
    setStatus(\"connecting\", retryCount === 0 ? \"연결 중…\" : `재연결 중… (${retryCount}/${MAX_RETRY})`);
    try {
        await client.join(\"APP_ID\", \"CHANNEL\", null, null);
        localTracks.audioTrack = await AgoraRTC.createMicrophoneAudioTrack({
            AEC: true, ANS: true, AGC: true,
        });
        await client.publish([localTracks.audioTrack]);
        client.enableAudioVolumeIndicator();

        client.on(\"volume-indicator\", (vs) => {
            vs.forEach((v) => {
                if (v.uid === 0 && !isMuted) {
                    document.getElementById(\"level-bar\").style.width = Math.min(v.level * 2, 100) + \"%\";
                }
            });
        });
        client.on(\"user-published\", async (u, m) => {
            try {
                await client.subscribe(u, m);
                if (m === \"audio\") u.audioTrack.play();
            } catch (e) { console.error(\"subscribe failed\", e); }
        });
        client.on(\"connection-state-change\", (cur, prev, reason) => {
            console.log(\"agora state:\", prev, \"->\", cur, reason);
            if (cur === \"CONNECTED\") {
                retryCount = 0;
                setStatus(\"connected\", \"🟢 연결됨\");
            } else if (cur === \"DISCONNECTED\") {
                setStatus(\"disconnected\", \"⚪ 연결 끊김\");
                if (retryCount < MAX_RETRY) {
                    retryCount += 1;
                    setTimeout(join, 1500 * retryCount);
                }
            } else if (cur === \"RECONNECTING\") {
                setStatus(\"connecting\", \"🔄 재연결 중…\");
            }
        });

        retryCount = 0;
        setStatus(\"connected\", \"🟢 연결됨\");
    } catch (e) {
        console.error(\"join failed\", e);
        if (retryCount < MAX_RETRY) {
            retryCount += 1;
            setStatus(\"connecting\", `재연결 시도… (${retryCount}/${MAX_RETRY})`);
            setTimeout(join, 1500 * retryCount);
        } else {
            setStatus(\"failed\", \"🔴 연결 실패 (새로고침)\");
        }
    }
}

function toggleMute() {
    if (!localTracks.audioTrack) return;
    isMuted = !isMuted;
    localTracks.audioTrack.setEnabled(!isMuted);
    const btn = document.getElementById(\"mute\");
    if (isMuted) {
        btn.innerText = \"🔇 마이크 : 음소거\";
        btn.classList.add(\"muted\");
    } else {
        btn.innerText = \"🎤 마이크 : 켜짐\";
        btn.classList.remove(\"muted\");
    }
}

join();
document.getElementById(\"mute\").onclick = toggleMute;
</script>
"""
    custom_html = (
        custom_html
        .replace("APP_ID", app_id)
        .replace("CHANNEL", channel)
        .replace("USER_LABEL", user_label)
    )
    components.html(custom_html, height=200)

@st.fragment(run_every="1s")
def sync_member_list(my_uid):
    with st.container(border=True):
        st.caption("실시간 보이스 연결 멤버")
        now = time.time()
        reporters, audience = [], []
        seen = set()
        for uid, info in shared_store["active_sessions"].items():
            if (now - info.get("last_seen", 0)) > 6:
                continue
            if not info.get("voice_connected"):
                continue
            label = info.get("label", "")
            if label in seen:
                continue
            seen.add(label)
            is_me = (uid == my_uid)
            display = label + (" (나)" if is_me else "")
            if info.get("role") == "reporter":
                reporters.append(display)
            else:
                audience.append(display)

        if not reporters and not audience:
            st.write("연결된 멤버 없음")
            return
        if reporters:
            st.markdown("**🎤 보고자**")
            for n in sorted(reporters):
                st.markdown(f"&nbsp;&nbsp;🟢 {n}")
        if audience:
            st.markdown("**👥 피보고자**")
            for n in sorted(audience):
                st.markdown(f"&nbsp;&nbsp;🟢 {n}")


# ==========================================
# 7. 사이드바
# ==========================================
with st.sidebar:
    st.title("AI Live Sync")
    st.caption(f"입장: **{my_label}**")
    if st.button("역할/이름 변경", use_container_width=True):
        for k in ("user_ready", "user_role", "user_name", "user_label"):
            st.session_state.pop(k, None)
        st.query_params.pop("role", None)
        st.query_params.pop("name", None)
        st.rerun()
    st.divider()

    voice_connect = st.toggle("음성 채널 접속", value=False, key="voice_active_toggle",
                              help="체크하면 Agora 보이스 채널에 접속합니다. 음소거는 별도 버튼에서 제어합니다.")

    if voice_connect:
        try:
            agora_id = st.secrets["AGORA_APP_ID"]
            agora_voice_system(agora_id, shared_store["voice_channel"], my_label)
        except Exception:
            st.warning("Agora ID 설정 필요")
    else:
        st.info("⚪ 음성 미접속 — 토글로 접속하세요")

    sync_member_list(st.session_state.uid)

    if is_reporter:
        st.divider()
        with st.expander("AI 자동 보고서 생성", expanded=False):
            try:
                ai_api_key = st.secrets["GEMINI_API_KEY"]
            except Exception:
                ai_api_key = ""
                st.warning("GEMINI_API_KEY 설정 필요")

            ai_text_input = st.text_area(
                "프롬프트 / 설명",
                placeholder="예) '청라 하늘대교 고률 점검 현장 보고서 3페이지. 구조물 교체 현황과 안전조치를 중점적으로 소개.'",
                height=120,
            )

            ai_file_input = st.file_uploader(
                "첨부 문서 (PDF / PPTX / DOCX / HWP / HWPX / TXT / CSV / MD)",
                type=["pdf", "pptx", "ppt", "docx", "hwp", "hwpx", "txt", "csv", "md"],
            )

            ai_photos = st.file_uploader(
                "📸 현장 사진 (jpg / png / webp · 다중 선택 · 모바일은 카메라 호출)",
                type=["jpg", "jpeg", "png", "webp"],
                accept_multiple_files=True,
                key="ai_photos_uploader",
                help="모바일에서 탭하면 '카메라'·'사진 선택' 옵션이 뜨고, 여러 장을 한 번에 올릴 수 있습니다.",
            )

            ai_camera = st.camera_input(
                "또는 즉석 촬영 (모바일 권장)",
                key="ai_camera_input",
                help="카메라가 에이는 경우 브라우저 주소창 자물쇠 → 카메라 허용을 확인하세요.",
            )

            if st.button("AI 보고서 생성", use_container_width=True, type="primary"):
                if not ai_api_key:
                    st.error("API Key 필요")
                else:
                    context = ai_text_input or ""
                    if ai_file_input:
                        doc_text = extract_text_from_upload(ai_file_input)
                        context += f"\n\n[첨부 문서: {ai_file_input.name}]\n{doc_text}"

                    # 사진 모아서 images 리스트 구성
                    photo_files = list(ai_photos or [])
                    if ai_camera is not None:
                        photo_files.append(ai_camera)
                    images = load_uploaded_images(photo_files)

                    if not context.strip() and not images:
                        st.error("프롬프트·문서·사진 중 최소 하나는 넣어주세요.")
                    else:
                        if not context.strip() and images:
                            context = "첨부된 현장 사진들을 분석해 현장 점검·진행현황 보고서를 작성해주세요."
                        req_pages = extract_requested_page_count(ai_text_input)
                        spinner_msg = "생성 중..."
                        if req_pages:
                            spinner_msg += f" ({req_pages}페이지)"
                        if images:
                            spinner_msg += f" · 사진 {len(images)}장 분석"
                        with st.spinner(spinner_msg):
                            ai_result = generate_json_from_ai(
                                ai_api_key, context,
                                requested_pages=req_pages,
                                images=images,
                            )
                        if "error" in ai_result:
                            st.error(f"생성 실패: {ai_result['error']}")
                        else:
                            ai_result = resolve_uploaded_tokens(ai_result, images)
                            shared_store["report_data"] = adapt_json_format(ai_result)
                            shared_store["current_page"] = 0
                            st.success(f"완료! 사진 {len(images)}장 관련되어 삽입됨." if images else "완료!")
                            time.sleep(1)
                            st.rerun()

        st.write("---")
        st.download_button(
            "표준 양식 다운로드",
            data=json.dumps(get_sample_json_guide(), indent=4, ensure_ascii=False),
            file_name="Report_Standard_Template.json",
            mime="application/json",
            use_container_width=True,
        )

        uploaded_file = st.file_uploader("JSON 수동 로드", type=["json"])
        if uploaded_file:
            if st.session_state.get("last_uploaded_id") != uploaded_file.file_id:
                shared_store["report_data"] = adapt_json_format(json.loads(uploaded_file.read().decode("utf-8")))
                st.session_state["last_uploaded_id"] = uploaded_file.file_id
                shared_store["current_page"] = 0

        if st.button("전체 데이터 초기화"):
            shared_store.update({"report_data": None, "current_page": 0, "chat_history": [], "active_sessions": {}})
            st.session_state.pop("last_uploaded_id", None)
            st.rerun()

        if shared_store["report_data"]:
            st.download_button(
                "최종 리포트 저장",
                data=json.dumps(shared_store["report_data"], indent=4, ensure_ascii=False),
                file_name="My_Final_Report.json",
                use_container_width=True,
            )

        edit_mode = st.toggle("편집 모드", value=False)
    else:
        edit_mode = False


# ==========================================
# 8. 메인 브리핑 엔진
# ==========================================
@st.fragment(run_every="1s")
def main_content_area(edit_enabled):
    shared_store["active_sessions"][st.session_state.uid] = {
        "label": my_label,
        "role": st.session_state.get("user_role", "audience"),
        "last_seen": time.time(),
        "voice_connected": st.session_state.get("voice_active_toggle", False),
    }

    with st.expander("실시간 채팅", expanded=False):
        c1, c2 = st.columns([4, 1])
        msg = c1.text_input("메시지", key="chat_in", label_visibility="collapsed")
        if c2.button("전송") and msg:
            shared_store["chat_history"].append(f"**{my_label}**: {msg}")
        chat_box = "".join([f"<div style='margin-bottom:6px;'>{m}</div>" for m in shared_store["chat_history"][-10:]])
        st.markdown(
            f"<div style='height:120px; overflow-y:auto; background:#f8f9fa; padding:12px; border-radius:10px; border:1px solid #dee2e6;'>{chat_box}</div>",
            unsafe_allow_html=True,
        )

    if shared_store["report_data"] is None:
        st.markdown(
            "<div style='text-align:center; padding:150px; color:#64748b;'><h2>좌측에서 AI 생성 또는 파일 로드</h2></div>",
            unsafe_allow_html=True,
        )
        if edit_enabled and st.button("새 보고서 시작"):
            shared_store["report_data"] = adapt_json_format({})
            st.rerun()
        return

    data = shared_store["report_data"]
    cp_idx = shared_store["current_page"]

    if edit_enabled:
        with st.expander("문서 공통 제목 설정", expanded=True):
            data["title"] = st.text_input("문서 제목", data.get("title", ""), key="global_title_input")
            dtc1, dtc2 = st.columns(2)
            data["title_fs"] = dtc1.slider("제목 글자 크기", 20, 120, int(data.get("title_fs", 55)))
            data["title_color"] = dtc2.color_picker("제목 색상", data.get("title_color", "#0f172a"))

    st.markdown(
        f'<h1 style="text-align:center; font-size:{data.get("title_fs", 55)}px; color:{data.get("title_color", "#0f172a")}; margin-bottom:10px;">{data.get("title", "")}</h1>',
        unsafe_allow_html=True,
    )
    st.markdown("<hr style='margin-top:0; margin-bottom:40px; border:0; border-top:1px solid #eee;'>", unsafe_allow_html=True)

    if cp_idx >= len(data["pages"]):
        shared_store["current_page"] = max(0, len(data["pages"]) - 1)
        st.rerun()

    p = data["pages"][cp_idx]

    if edit_enabled:
        st.write("---")
        pc1, pc2 = st.columns([1, 5])
        if pc1.button("페이지 추가"):
            data["pages"].insert(cp_idx + 1, create_empty_page())
            shared_store["current_page"] += 1
            st.rerun()
        if pc2.button("페이지 삭제") and len(data["pages"]) > 1:
            data["pages"].pop(cp_idx)
            shared_store["current_page"] = max(0, cp_idx - 1)
            st.rerun()

    if is_reporter:
        tabs = {i: f"P{i+1}. {pg.get('tab', '')}" for i, pg in enumerate(data["pages"])}
        shared_store["current_page"] = st.radio(
            "이동", list(tabs.keys()),
            index=shared_store["current_page"],
            format_func=lambda x: tabs[x],
            horizontal=True,
        )

    if edit_enabled:
        p["tab"] = st.text_input("탭 이름", p.get("tab", ""), key=f"tab_edit{cp_idx}")
        with st.expander("소제목 설정"):
            p["header"] = st.text_input("소제목", p.get("header", ""), key=f"ph_input{cp_idx}")
            hc1, hc2 = st.columns(2)
            p["header_fs"] = hc1.slider("크기", 10, 150, int(p.get("header_fs", 35)), key=f"phfs{cp_idx}")
            p["header_color"] = hc2.color_picker("색상", p.get("header_color", "#475569"), key=f"phc{cp_idx}")

    st.markdown(
        f'<h2 style="text-align:center; font-size:{p.get("header_fs", 35)}px; color:{p.get("header_color", "#475569")}; margin-bottom:30px;">{p.get("header", "")}</h2>',
        unsafe_allow_html=True,
    )

    sections = p.setdefault("sections", [])
    if edit_enabled and st.button("섹션 추가", key=f"add_sec_btn{cp_idx}"):
        sections.append({
            "title": "새 섹션", "title_fs": 32, "title_color": "#1a1c1e", "col_ratio": 1.5,
            "lines": [{"text": "내용", "size": 22, "color": "#1e293b"}],
            "main_image": None, "full_width": True, "image_query": "",
            "chart_type": "Bar", "chart_data": "", "side_items": [],
        })
        st.rerun()

    for s_idx, sec in enumerate(sections):
        with st.container(border=True):
            if edit_enabled:
                sc1, sc2, sc3, sc4, sc5 = st.columns([2.5, 0.8, 0.8, 1.2, 0.5])
                sec["title"] = sc1.text_input("섹션 제목", sec.get("title", ""), key=f"st_t{cp_idx}{s_idx}")
                sec["title_fs"] = sc2.number_input("크기", 10, 80, int(sec.get("title_fs", 32)), key=f"st_fs{cp_idx}{s_idx}")
                sec["title_color"] = sc3.color_picker("색상", sec.get("title_color", "#1a1c1e"), key=f"st_c{cp_idx}{s_idx}")
                sec["col_ratio"] = sc4.slider("비율", 1.0, 4.0, float(sec.get("col_ratio", 1.5)), 0.1, key=f"st_r{cp_idx}{s_idx}")
                if sc5.button("X", key=f"st_del{cp_idx}{s_idx}"):
                    sections.pop(s_idx)
                    st.rerun()

            st.markdown(
                f"<h3 style='font-size:{sec.get('title_fs', 32)}px; color:{sec.get('title_color', '#1a1c1e')}; margin-bottom:20px; padding-bottom:12px; border-bottom:2px solid #f8f9fa;'>{sec.get('title')}</h3>",
                unsafe_allow_html=True,
            )

            col_main, col_side = st.columns([sec.get("col_ratio", 1.5), 1], gap="medium")

            with col_main:
                if edit_enabled:
                    with st.expander("이미지 관리"):
                        current_img = sec.get("main_image", "") or ""
                        display_img_url = "" if (isinstance(current_img, str) and current_img.startswith("data:image")) else current_img
                        new_img_url = st.text_input("URL", value=display_img_url, key=f"simg_url{cp_idx}{s_idx}")
                        sec["image_query"] = st.text_input("자동 검색어 (프로젝트 고유명사+핵심어)", value=sec.get("image_query", ""), key=f"simg_q{cp_idx}{s_idx}")
                        img_f = st.file_uploader("업로드", type=["png", "jpg", "jpeg", "gif"], key=f"simg_f{cp_idx}{s_idx}")
                        if img_f:
                            sec["main_image"] = f"data:image/png;base64,{base64.b64encode(img_f.getvalue()).decode()}"
                        elif new_img_url != display_img_url:
                            sec["main_image"] = new_img_url
                        sec["full_width"] = st.toggle("너비 채우기", value=sec.get("full_width", True), key=f"fw{cp_idx}{s_idx}")
                        if not sec["full_width"]:
                            sec["img_width"] = st.slider("너비", 100, 1200, int(sec.get("img_width", 750)), key=f"sw{cp_idx}{s_idx}")
                        if st.button("그림 삭제", key=f"simg_del{cp_idx}{s_idx}"):
                            sec["main_image"] = None
                            st.rerun()

                if (not sec.get("main_image")) and sec.get("image_query"):
                    sec["main_image"] = get_auto_image_url(sec.get("image_query"))

                if sec.get("main_image"):
                    final_src = render_image_src(sec["main_image"])
                    style = "width:100%;" if sec.get("full_width", True) else f"width:{sec.get('img_width', 750)}px; max-width:100%;"
                    st.markdown(
                        f'<div style="text-align:center;"><img src="{final_src}" onerror="this.onerror=null; this.src={SQ}{ERROR_IMG}{SQ};" style="{style} border-radius:12px; margin-bottom:20px; box-shadow:0 4px 12px rgba(0,0,0,0.05);" /></div>',
                        unsafe_allow_html=True,
                    )

                if edit_enabled:
                    with st.expander("차트 관리"):
                        sec["chart_type"] = st.selectbox(
                            "종류", ["Bar", "Line", "Area"],
                            index=["Bar", "Line", "Area"].index(sec.get("chart_type", "Bar")) if sec.get("chart_type") in ["Bar", "Line", "Area"] else 0,
                            key=f"ch_t{cp_idx}{s_idx}",
                        )
                        sec["chart_data"] = st.text_area("데이터 (항목, 숫자)", value=sec.get("chart_data", ""), key=f"ch_d{cp_idx}{s_idx}")

                if sec.get("chart_data"):
                    try:
                        raw_chart = sec["chart_data"].replace("\\n", "\n")
                        lines_data = [line.strip() for line in raw_chart.split("\n") if "," in line]
                        if lines_data:
                            data_dict = {}
                            for line in lines_data:
                                k, v = line.split(",", 1)
                                data_dict[k.strip()] = float(v.replace(",", "").strip())
                            if data_dict:
                                df = pd.DataFrame(list(data_dict.values()), index=list(data_dict.keys()), columns=["수치"])
                                ctype = sec.get("chart_type", "Bar")
                                if ctype == "Line":
                                    st.line_chart(df, use_container_width=True)
                                elif ctype == "Area":
                                    st.area_chart(df, use_container_width=True)
                                else:
                                    st.bar_chart(df, use_container_width=True)
                    except Exception:
                        st.warning("차트 데이터 형식 오류")

                sec.setdefault("lines", [])
                if edit_enabled:
                    st.caption("본문 편집")
                    new_lines = []
                    for l_idx, line in enumerate(sec["lines"]):
                        lc1, lc2, lc3, lc4 = st.columns([5, 1.5, 1.5, 0.5])
                        l_t = lc1.text_input("T", line["text"], key=f"lt_t{cp_idx}{s_idx}{l_idx}", label_visibility="collapsed")
                        l_s = lc2.number_input("S", 10, 100, int(line["size"]), key=f"lt_s_{cp_idx}{s_idx}{l_idx}")
                        l_c = lc3.color_picker("C", line["color"], key=f"lt_c_{cp_idx}{s_idx}{l_idx}")
                        if not lc4.button("X", key=f"lt_del_{cp_idx}{s_idx}{l_idx}"):
                            new_lines.append({"text": l_t, "size": l_s, "color": l_c})
                    sec["lines"] = new_lines
                    if st.button("줄 추가", key=f"lt_add_{cp_idx}{s_idx}"):
                        sec["lines"].append({"text": "새 문구", "size": 22, "color": "#1e293b"})
                        st.rerun()
                else:
                    for line in sec.get("lines", []):
                        st.markdown(
                            f'<p class="text-line" style="font-size:{line["size"]}px; color:{line["color"]}; font-weight:bold;">{line["text"]}</p>',
                            unsafe_allow_html=True,
                        )

            with col_side:
                sec.setdefault("side_items", [])
                if edit_enabled:
                    sca, scb = st.columns(2)
                    if sca.button("지표 추가", key=f"si_add_m{cp_idx}{s_idx}"):
                        sec["side_items"].append({"type": "metric", "label": "항목", "value": "0", "color": "#007bff", "label_fs": 14, "label_color": "#64748b", "value_fs": 28})
                        st.rerun()
                    if scb.button("그림 추가", key=f"si_add_i{cp_idx}{s_idx}"):
                        sec["side_items"].append({"type": "image", "src": None, "width": 350, "image_query": ""})
                        st.rerun()

                for i_idx, item in enumerate(sec["side_items"]):
                    if edit_enabled:
                        with st.expander(f"{item.get('label', '아이템')} 편집", expanded=True):
                            if item["type"] == "metric":
                                item["label"] = st.text_input("라벨", item.get("label"), key=f"si_l{cp_idx}{s_idx}{i_idx}")
                                item["value"] = st.text_area("내용", item.get("value"), height=120, key=f"si_v_{cp_idx}{s_idx}{i_idx}")
                                ic3, ic4 = st.columns(2)
                                item["label_fs"] = ic3.number_input("라벨크기", 10, 60, int(item.get("label_fs", 14)), key=f"si_lfs_{cp_idx}{s_idx}{i_idx}")
                                item["label_color"] = ic4.color_picker("라벨색상", item.get("label_color", "#64748b"), key=f"si_lc_{cp_idx}{s_idx}{i_idx}")
                                ic5, ic6 = st.columns(2)
                                item["value_fs"] = ic5.number_input("내용크기", 10, 100, int(item.get("value_fs", 28)), key=f"si_vfs_{cp_idx}{s_idx}{i_idx}")
                                item["color"] = ic6.color_picker("내용색상", item.get("color", "#007bff"), key=f"si_vc_{cp_idx}{s_idx}{i_idx}")
                            elif item["type"] == "image":
                                current_side = item.get("src", "") or ""
                                display_side = "" if (isinstance(current_side, str) and current_side.startswith("data:image")) else current_side
                                new_side_url = st.text_input("URL", value=display_side, key=f"si_url_{cp_idx}{s_idx}{i_idx}")
                                item["image_query"] = st.text_input("검색어(프로젝트 고유명사+핵심어)", value=item.get("image_query", ""), key=f"si_q_{cp_idx}{s_idx}{i_idx}")
                                siu = st.file_uploader("업로드", type=["png", "jpg", "jpeg", "gif"], key=f"si_iu_{cp_idx}{s_idx}{i_idx}")
                                if siu:
                                    item["src"] = f"data:image/png;base64,{base64.b64encode(siu.getvalue()).decode()}"
                                elif new_side_url != display_side:
                                    item["src"] = new_side_url
                                item["width"] = st.slider("너비", 100, 500, int(item.get("width", 350)), key=f"si_iw_{cp_idx}{s_idx}{i_idx}")
                            if st.button("삭제", key=f"si_del_{cp_idx}{s_idx}{i_idx}"):
                                sec["side_items"].pop(i_idx)
                                st.rerun()

                    if item["type"] == "metric":
                        raw_val = item.get("value", "")
                        if isinstance(raw_val, str):
                            raw_val = raw_val.replace("/n", "\n")
                        fv = (raw_val or "").replace("\n", "<br>")
                        st.markdown(
                            f'<div class="side-slot-card">'
                            f'<div style="font-size:{item.get("label_fs", 14)}px; color:{item.get("label_color", "#64748b")}; margin-bottom:8px;">{item.get("label", "")}</div>'
                            f'<div style="font-size:{item.get("value_fs", 28)}px; font-weight:bold; color:{item.get("color", "#007bff")}; line-height:1.5;">{fv}</div>'
                            f'</div>',
                            unsafe_allow_html=True,
                        )
                    elif item["type"] == "image":
                        if (not item.get("src")) and item.get("image_query"):
                            item["src"] = get_auto_image_url(item.get("image_query"), w=900, h=700)
                        if item.get("src"):
                            final_side_src = render_image_src(item["src"])
                            st.markdown(
                                f'<div class="side-slot-card">'
                                f'<img src="{final_side_src}" onerror="this.onerror=null; this.src={SQ}{ERROR_IMG}{SQ};" style="width:{item.get("width", 350)}px; max-width:100%; border-radius:12px; box-shadow:0 4px 12px rgba(0,0,0,0.08);" />'
                                f'</div>',
                                unsafe_allow_html=True,
                            )


# 실행
main_content_area(edit_mode)
